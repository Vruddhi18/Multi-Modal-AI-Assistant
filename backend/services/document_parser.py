import PyPDF2
from pptx import Presentation
from docx import Document
import openpyxl
import csv
import io
from PIL import Image
import pytesseract

IMAGE_EXTS = {'png', 'jpg', 'jpeg', 'bmp', 'tiff', 'webp'}

def parse_document(file_bytes: bytes, ext: str) -> str:
    text = ""
    try:
        if ext == 'pdf':
            pdf_file = io.BytesIO(file_bytes)
            reader = PyPDF2.PdfReader(pdf_file)
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"

        elif ext == 'pptx':
            pptx_file = io.BytesIO(file_bytes)
            prs = Presentation(pptx_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        elif ext == 'docx':
            docx_file = io.BytesIO(file_bytes)
            doc = Document(docx_file)
            for para in doc.paragraphs:
                if para.text:
                    text += para.text + "\n"
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text:
                            text += cell.text + "\n"

        elif ext == 'xlsx':
            xlsx_file = io.BytesIO(file_bytes)
            wb = openpyxl.load_workbook(xlsx_file, data_only=True)
            for sheet in wb.worksheets:
                text += f"Sheet: {sheet.title}\n"
                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(cell) for cell in row if cell is not None]
                    if row_values:
                        text += ", ".join(row_values) + "\n"

        elif ext == 'csv':
            decoded = file_bytes.decode('utf-8', errors='ignore')
            reader_csv = csv.reader(io.StringIO(decoded))
            for row in reader_csv:
                if row:
                    text += ", ".join(row) + "\n"

        elif ext in IMAGE_EXTS:
            image = Image.open(io.BytesIO(file_bytes))
            text = pytesseract.image_to_string(image)

        elif ext == 'txt':
            text = file_bytes.decode('utf-8')

    except Exception as e:
        print(f"Error parsing document: {e}")
        raise e

    return text